"""File-format detection and text extraction for Athena.

Resolves a URL or local file to a canonical source_type + extracted
prose text, so the ingest pipeline can build a wiki page without
guessing from URL suffixes alone.

Two public entry points:

    detect_type(url, timeout=10) -> {
        'source_type': 'paper'|'doc'|'spreadsheet'|'slides'|'webpage'|'image'|...,
        'ext':         'pdf'|'docx'|'xlsx'|'pptx'|'epub'|'md'|'txt'|'csv'|...,
        'mime':        'application/pdf'|...,
        'filename':    '<server-suggested or URL-basename>',
    }

    extract_text(filepath, ext) -> {
        'title':   '<doc title or None>',
        'authors': '<comma-joined authors or None>',
        'text':    '<plain prose, may be long>',
    }

Design notes:
    - Office 2007+ formats (.docx/.xlsx/.pptx) are ZIP+XML — handled
      with Python stdlib (zipfile + xml.etree). No pip dependencies.
    - PDF extraction uses `pdftotext` (from poppler) when available.
      If poppler isn't installed, we fall back to metadata-only via
      `pdfinfo`, then finally to filename-based title.
    - EPUB is also ZIP+XHTML — handled the same way as Office.
    - detect_type uses a HEAD request to read the server's own
      Content-Type + Content-Disposition headers. URL-suffix matching
      is a fallback for servers that don't return these headers.
"""

import os
import re
import zipfile
import subprocess
import urllib.parse
import urllib.request
import urllib.error
import xml.etree.ElementTree as ET


# MIME-to-source-type mapping. source_type values match Athena's
# existing taxonomy (wiki/format/{papers,repos,webpages,videos,images})
# plus three new folders created on demand: docs, spreadsheets, slides.
_MIME_MAP = {
    'application/pdf': ('paper', 'pdf'),
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': ('doc', 'docx'),
    'application/msword': ('doc', 'doc'),
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': ('spreadsheet', 'xlsx'),
    'application/vnd.ms-excel': ('spreadsheet', 'xls'),
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': ('slides', 'pptx'),
    'application/vnd.ms-powerpoint': ('slides', 'ppt'),
    'application/epub+zip': ('book', 'epub'),
    'text/markdown': ('webpage', 'md'),
    'text/plain': ('webpage', 'txt'),
    'text/csv': ('spreadsheet', 'csv'),
    'text/html': ('webpage', 'html'),
    'image/png': ('image', 'png'),
    'image/jpeg': ('image', 'jpg'),
    'image/gif': ('image', 'gif'),
    'image/webp': ('image', 'webp'),
}

# URL-suffix fallback (checked if HEAD fails or returns a generic type).
_EXT_MAP = {
    'pdf': ('paper', 'pdf'),
    'docx': ('doc', 'docx'),
    'doc': ('doc', 'doc'),
    'xlsx': ('spreadsheet', 'xlsx'),
    'xls': ('spreadsheet', 'xls'),
    'pptx': ('slides', 'pptx'),
    'ppt': ('slides', 'ppt'),
    'epub': ('book', 'epub'),
    'md': ('webpage', 'md'),
    'txt': ('webpage', 'txt'),
    'csv': ('spreadsheet', 'csv'),
    'png': ('image', 'png'),
    'jpg': ('image', 'jpg'),
    'jpeg': ('image', 'jpg'),
}


# ── Content-Type detection ──────────────────────────────────────

def detect_type(url, timeout=10):
    """Resolve a URL to {source_type, ext, mime, filename}.

    Tries HEAD first (reads Content-Type + Content-Disposition). Falls
    back to URL-suffix parsing if HEAD fails or returns a generic type.
    Never raises — returns 'webpage'/'html' as the catch-all default.
    """
    out = {
        'source_type': 'webpage',
        'ext': 'html',
        'mime': None,
        'filename': None,
    }

    # Platform-specific URL patterns short-circuit all the detection —
    # we want GitHub/arxiv/youtube to go through their bespoke flows.
    if re.search(r'github\.com/[^/]+/[^/]+', url, re.IGNORECASE):
        out.update({'source_type': 'repo', 'ext': None, 'mime': None})
        return out
    if re.search(r'(youtube\.com|youtu\.be)', url, re.IGNORECASE):
        out.update({'source_type': 'video', 'ext': None, 'mime': None})
        return out
    if re.search(r'arxiv\.org/abs/', url, re.IGNORECASE):
        out.update({'source_type': 'paper', 'ext': 'pdf', 'mime': 'application/pdf'})
        return out
    # Google Drive `/file/d/<id>/view` URLs are PDFs (or other files)
    # served through a viewer UI — HEAD on the view URL returns HTML,
    # so Content-Type detection would miss them. We assume PDF since
    # that's the realistic case in Athena's ingest workflow; a non-PDF
    # Drive file will still get caught by magic-byte validation during
    # download (ingest-file discards non-PDF content).
    if re.search(r'drive\.google\.com/file/d/[^/]+', url, re.IGNORECASE):
        out.update({'source_type': 'paper', 'ext': 'pdf', 'mime': 'application/pdf'})
        return out

    # Try HEAD request. Some servers don't support HEAD, others respond
    # with wrong headers — we accept whatever we can parse.
    try:
        req = urllib.request.Request(
            url, method='HEAD',
            headers={'User-Agent': 'Athena/1.0 (content detector)'},
        )
        with urllib.request.urlopen(req, timeout=timeout) as resp:
            mime = (resp.headers.get('Content-Type') or '').split(';')[0].strip().lower()
            disp = resp.headers.get('Content-Disposition') or ''
            if mime:
                out['mime'] = mime
                if mime in _MIME_MAP:
                    out['source_type'], out['ext'] = _MIME_MAP[mime]
            # Parse filename from Content-Disposition if present — this
            # is the server's own suggested filename, usually more
            # meaningful than the URL path.
            fn_match = re.search(r'filename\*?=(?:[^\'"]*\'\')?"?([^";]+)"?', disp)
            if fn_match:
                out['filename'] = urllib.parse.unquote(fn_match.group(1).strip())
    except (urllib.error.URLError, urllib.error.HTTPError, TimeoutError, OSError):
        pass  # HEAD failed — fall through to URL-suffix detection

    # URL-suffix fallback: runs if HEAD didn't produce a known mapping.
    # The "generic" Content-Types ("application/octet-stream", etc.) also
    # fall here — we trust the URL over a vague server response.
    if out['source_type'] == 'webpage' and out['ext'] == 'html':
        path = urllib.parse.urlparse(url).path or ''
        ext_match = re.search(r'\.([a-z0-9]+)(?:$|\?)', path, re.IGNORECASE)
        if ext_match:
            ext = ext_match.group(1).lower()
            if ext in _EXT_MAP:
                out['source_type'], out['ext'] = _EXT_MAP[ext]
                if not out['mime']:
                    # Best-guess MIME from ext — not authoritative, but
                    # useful if downstream code keys on mime.
                    rev = {v: k for k, v in _MIME_MAP.items()}
                    out['mime'] = rev.get((out['source_type'], out['ext']))

    return out


# ── Text extraction ─────────────────────────────────────────────

def extract_text(filepath, ext):
    """Return {title, authors, text} for a downloaded file.

    Empty strings rather than None so callers can concatenate freely.
    Failures (missing file, binary garbage, missing helper tool) return
    empty fields with no exception — ingest proceeds with whatever
    metadata it has.
    """
    out = {'title': '', 'authors': '', 'text': ''}
    if not filepath or not os.path.exists(filepath):
        return out

    ext = (ext or '').lower()
    try:
        if ext == 'pdf':
            out.update(_extract_pdf(filepath))
        elif ext == 'docx':
            out.update(_extract_docx(filepath))
        elif ext == 'xlsx':
            out.update(_extract_xlsx(filepath))
        elif ext == 'pptx':
            out.update(_extract_pptx(filepath))
        elif ext == 'epub':
            out.update(_extract_epub(filepath))
        elif ext in ('md', 'txt', 'csv'):
            with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
                out['text'] = f.read()
    except (OSError, zipfile.BadZipFile, ET.ParseError):
        # Stay silent — missing text is recoverable; a noisy stderr
        # pollutes the JSON that wiki_page.py emits to stdout.
        pass

    return out


def _run_tool(cmd, timeout=30):
    """Subprocess helper that returns stdout or empty string on failure."""
    try:
        result = subprocess.run(
            cmd, capture_output=True, text=True, timeout=timeout, check=False,
        )
        return result.stdout if result.returncode == 0 else ''
    except (subprocess.TimeoutExpired, FileNotFoundError, OSError):
        return ''


def _extract_pdf(filepath):
    """Tiered PDF extraction:
      1) pymupdf4llm — pip install, produces structured markdown
         (preserves headings, lists, code blocks; handles tables better
         than pdftotext).
      2) pdftotext — poppler CLI, plain-text fallback, always available
         after `brew install poppler`.
      3) empty — raw file gets filename-derived title from the caller.

    Metadata (title, authors) always comes from pdfinfo since even the
    best body extractor doesn't expose XMP/Info cleanly.

    Returns {title, authors, text, tier, pages}. `pages` is a list of
    {"page": <1-indexed>, "text": <markdown>} when the structured tier
    ran (parallel to the provider's segments); `tier` records which
    extractor produced the body ('pymupdf4llm', 'pdftotext', or '').

    The page number is read from the chunk's own metadata so locators
    track the PDF's real page identity (including non-contiguous page-range
    extracts), not a sequential counter. Precedence: the installed
    pymupdf4llm exposes a 1-indexed `page_number` key (verified against
    1.27.2.3); a legacy 0-indexed `page` key is supported as a fallback
    (+1 applied); absent both, a 1-based sequential counter is used."""
    out = {'title': '', 'authors': '', 'text': '', 'tier': '', 'pages': []}
    info = _run_tool(['pdfinfo', filepath], timeout=10)
    for line in info.splitlines():
        if line.startswith('Title:') and not out['title']:
            out['title'] = line.split(':', 1)[1].strip()
        elif line.startswith('Author:') and not out['authors']:
            out['authors'] = line.split(':', 1)[1].strip()

    # Try pymupdf4llm page chunks first. Imported lazily because it's a
    # pip dep users may not have; absence shouldn't break ingest.
    page_chunks = _pdf_pages_via_pymupdf4llm(filepath)
    if page_chunks:
        out['tier'] = 'pymupdf4llm'
        pages = []
        for chunk in page_chunks:
            meta = chunk.get('metadata') or {}
            if 'page_number' in meta:        # real pymupdf4llm key, 1-indexed
                page_no = meta['page_number']
            elif 'page' in meta:             # 0-indexed (older versions / test fixtures)
                page_no = meta['page'] + 1
            else:
                page_no = len(pages) + 1     # last-resort sequential fallback
            pages.append({'page': page_no, 'text': (chunk.get('text') or '').strip()})
        out['pages'] = pages
        out['text'] = ('\n\n'.join(p['text'] for p in pages))[:50000]
        return out

    # Fallback: pdftotext default mode. `-nopgbrk` drops form-feed
    # page breaks. Reflow merges soft-wrapped lines into paragraphs.
    raw = _run_tool(['pdftotext', '-nopgbrk', filepath, '-'], timeout=60)
    out['tier'] = 'pdftotext' if raw else ''
    out['text'] = _reflow_paragraphs(raw)[:50000]
    return out


def _pdf_pages_via_pymupdf4llm(filepath):
    """Per-page markdown chunks via pymupdf4llm. Returns a list of
    {"text", "metadata": {...}} dicts, or [] on any failure. The metadata
    dict carries a 1-indexed `page_number` key in the installed version
    (verified against pymupdf4llm 1.27.2.3); there is no `page` key.
    `_extract_pdf` reads `page_number` to derive accurate page locators.

    pymupdf4llm prints layout/version advisories to stdout — we redirect
    both stdout and stderr during the call so these don't pollute the
    JSON our caller writes to its stdout. Missing-module isn't an error
    condition — it's a graceful tier."""
    import contextlib
    import io
    try:
        import pymupdf4llm  # noqa: F401
    except ImportError:
        return []
    try:
        with contextlib.redirect_stdout(io.StringIO()), \
             contextlib.redirect_stderr(io.StringIO()):
            return pymupdf4llm.to_markdown(
                filepath, page_chunks=True, show_progress=False
            ) or []
    except Exception:
        return []


def _reflow_paragraphs(text):
    """Merge soft-wrapped lines inside paragraphs. pdftotext wraps at
    the PDF's column width, which means a single sentence often spans
    3-4 lines. Join lines that are clearly mid-paragraph (line ends in
    a letter, comma, or semicolon) and keep blank lines intact as
    paragraph separators.

    Heuristic, not perfect — lists and headings can get mis-joined.
    The downstream LLM handles residual noise fine; this is purely a
    readability improvement for the raw file."""
    out_lines = []
    buf = []
    for line in text.splitlines():
        stripped = line.strip()
        if not stripped:
            if buf:
                out_lines.append(' '.join(buf))
                buf = []
            out_lines.append('')
            continue
        buf.append(stripped)
        # End-of-paragraph markers: ., !, ?, : followed by nothing.
        # Everything else (including commas, semicolons, mid-word splits)
        # is a soft wrap we'll join.
        if re.search(r'[.!?:"\)]\s*$', stripped):
            out_lines.append(' '.join(buf))
            buf = []
    if buf:
        out_lines.append(' '.join(buf))
    # Collapse runs of blank lines to single blanks for cleaner output.
    result = '\n'.join(out_lines)
    return re.sub(r'\n{3,}', '\n\n', result)


def _zip_text_from(zip_path, entry_paths, text_tag):
    """Shared core for Office/EPUB extraction: open the ZIP, read a list
    of entry paths (some may be missing), concatenate the text inside
    the given XML tag. Returns a single flat string of paragraphs."""
    chunks = []
    with zipfile.ZipFile(zip_path) as zf:
        existing = set(zf.namelist())
        for entry in entry_paths:
            if entry not in existing:
                continue
            try:
                raw = zf.read(entry)
                root = ET.fromstring(raw)
            except (ET.ParseError, KeyError):
                continue
            # Collect <text_tag> text across all namespaces — the
            # default-namespace declarations in Office XML mean the
            # literal tag name carries a "{uri}tag" prefix at parse
            # time, so we match by local name.
            for el in root.iter():
                tag = el.tag.rsplit('}', 1)[-1] if '}' in el.tag else el.tag
                if tag == text_tag and el.text:
                    chunks.append(el.text)
    return '\n'.join(chunks)


def _pandoc_to_markdown(filepath, input_format=None):
    """Run pandoc to convert a file to GitHub-flavored markdown. Returns
    '' when pandoc isn't installed or conversion fails — callers fall
    back to the stdlib extractors. GFM is chosen over plain markdown
    because Obsidian renders it identically and the table syntax is
    easier to read in the raw file."""
    cmd = ['pandoc']
    if input_format:
        cmd += ['--from', input_format]
    cmd += ['--to', 'gfm', '--wrap=none', filepath]
    return _run_tool(cmd, timeout=60)


_ENTRY_NUM = re.compile(r'(\d+)\.xml$')


def _entry_index(name):
    """Numeric index from a `...<N>.xml` entry name (e.g. slide10.xml → 10).

    Lexicographic sorting would order `slide10.xml` before `slide2.xml`,
    mis-numbering locators for documents with 10+ slides/sheets (R5). Sort
    by this numeric key instead. Returns 0 for names with no trailing index."""
    m = _ENTRY_NUM.search(name)
    return int(m.group(1)) if m else 0


def _pptx_units(filepath):
    """Per-slide text from a pptx, 1-indexed by slide order (stdlib path)."""
    units = []
    try:
        with zipfile.ZipFile(filepath) as zf:
            slide_entries = sorted(
                (n for n in zf.namelist()
                 if n.startswith("ppt/slides/slide") and n.endswith(".xml")),
                key=_entry_index,
            )
    except (zipfile.BadZipFile, OSError):
        return units
    for i, entry in enumerate(slide_entries, start=1):
        text = _zip_text_from(filepath, [entry], "t").strip()
        units.append({"slide": i, "text": text})
    return units


def _xlsx_units(filepath):
    """Per-sheet text from an xlsx. Sheet name from xl/workbook.xml in
    declaration order, paired with xl/worksheets/sheetN.xml in name-sorted
    order. This pairing is correct for the common case where sheet files are
    created in declaration order; it degrades to positional labels if not."""
    units = []
    try:
        with zipfile.ZipFile(filepath) as zf:
            names = _xlsx_sheet_names(filepath)  # ordered list of sheet names
            sheet_entries = sorted(
                (n for n in zf.namelist()
                 if n.startswith("xl/worksheets/sheet") and n.endswith(".xml")),
                key=_entry_index,
            )
    except (zipfile.BadZipFile, OSError):
        return units
    for i, entry in enumerate(sheet_entries):
        sheet_name = names[i] if i < len(names) else f"Sheet{i+1}"
        text = _zip_text_from(filepath, [entry], "t").strip()
        units.append({"sheet": sheet_name, "text": text})
    return units


def _xlsx_sheet_names(filepath):
    """Ordered sheet names from xl/workbook.xml; [] on failure."""
    names = []
    try:
        with zipfile.ZipFile(filepath) as zf:
            if "xl/workbook.xml" not in zf.namelist():
                return names
            root = ET.fromstring(zf.read("xl/workbook.xml"))
    except (zipfile.BadZipFile, ET.ParseError, OSError, KeyError):
        return names
    for el in root.iter():
        tag = el.tag.rsplit("}", 1)[-1] if "}" in el.tag else el.tag
        if tag == "sheet":
            name = el.get("name")
            if name:
                names.append(name)
    return names


# ── pure-pip structured tier ([office] extra: openpyxl / python-pptx / python-docx) ──
# These produce structured Markdown WITHOUT a system binary (pandoc) or the heavy
# Docling models, filling the pip-native gap between Docling and the flat zip walk.
# Each returns '' / ([], '') when its library is absent or the file can't be parsed,
# so the caller falls through to the next tier.


def _rows_to_gfm(rows):
    """List-of-rows (each a list of cell strings) → a GitHub-flavored Markdown
    table. The first row is the header. Trims fully-empty trailing rows/columns;
    returns '' when there's no data. Pipes/newlines in cells are escaped/flattened."""
    rows = [list(r) for r in rows]
    while rows and not any((c or '').strip() for c in rows[-1]):
        rows.pop()
    if not rows:
        return ''
    ncols = max((len(r) for r in rows), default=0)
    while ncols > 0 and all(len(r) < ncols or not (r[ncols - 1] or '').strip() for r in rows):
        ncols -= 1
    if ncols == 0:
        return ''

    def cell(r, i):
        return (r[i] if i < len(r) else '').replace('|', r'\|').replace('\n', ' ').strip()

    out = ['| ' + ' | '.join(cell(rows[0], i) for i in range(ncols)) + ' |']
    out.append('| ' + ' | '.join('---' for _ in range(ncols)) + ' |')
    for r in rows[1:]:
        out.append('| ' + ' | '.join(cell(r, i) for i in range(ncols)) + ' |')
    return '\n'.join(out)


def _docx_via_python_docx(filepath):
    """Structured docx → Markdown via python-docx. '' if the lib is missing or the
    file can't be parsed. Heading paragraphs → `#`×level; tables → GFM; Title → H1."""
    try:
        import docx
        from docx.oxml.ns import qn
        from docx.table import Table
        from docx.text.paragraph import Paragraph
    except ImportError:
        return ''
    try:
        document = docx.Document(filepath)
    except Exception:
        return ''
    lines = []
    for child in document.element.body.iterchildren():
        if child.tag == qn('w:p'):
            para = Paragraph(child, document)
            text = para.text.strip()
            if not text:
                continue
            style = (para.style.name if para.style else '') or ''
            if style.startswith('Heading'):
                digits = ''.join(ch for ch in style if ch.isdigit())
                lines.append('#' * (min(int(digits), 6) if digits else 1) + ' ' + text)
            elif style == 'Title':
                lines.append('# ' + text)
            else:
                lines.append(text)
        elif child.tag == qn('w:tbl'):
            md = _rows_to_gfm([[c.text.strip() for c in row.cells] for row in Table(child, document).rows])
            if md:
                lines.append(md)
    return '\n\n'.join(lines).strip()


def _xlsx_via_openpyxl(filepath):
    """Structured xlsx → (markdown, per-sheet units) via openpyxl. Each sheet is a
    `## <name>` section with a GFM table. ('', []) if the lib is missing/parse fails."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        return '', []
    try:
        wb = load_workbook(filepath, read_only=True, data_only=True)
    except Exception:
        return '', []
    sections, units = [], []
    try:
        for ws in wb.worksheets:
            rows = [['' if c is None else str(c) for c in row]
                    for row in ws.iter_rows(values_only=True)]
            table_md = _rows_to_gfm(rows)
            units.append({'sheet': ws.title, 'text': table_md})
            sections.append(f'## {ws.title}\n\n{table_md}' if table_md else f'## {ws.title}')
    finally:
        wb.close()
    return '\n\n'.join(sections).strip(), units


def _pptx_via_python_pptx(filepath):
    """Structured pptx → (markdown, per-slide units) via python-pptx. Each slide is a
    `## Slide N` section with its text frames + tables. ('', []) if missing/parse fails."""
    try:
        from pptx import Presentation
    except ImportError:
        return '', []
    try:
        prs = Presentation(filepath)
    except Exception:
        return '', []
    sections, units = [], []
    for i, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if getattr(shape, 'has_table', False) and shape.has_table:
                md = _rows_to_gfm([[c.text.strip() for c in row.cells] for row in shape.table.rows])
                if md:
                    parts.append(md)
            elif getattr(shape, 'has_text_frame', False) and shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if txt:
                    parts.append(txt)
        slide_md = '\n\n'.join(parts).strip()
        units.append({'slide': i, 'text': slide_md})
        sections.append(f'## Slide {i}\n\n{slide_md}' if slide_md else f'## Slide {i}')
    return '\n\n'.join(sections).strip(), units


def _extract_docx(filepath):
    """Word 2007+: python-docx (pure-pip [office] extra) is the primary structured
    tier; pandoc is next; stdlib XML walking is the flat fallback. Title from
    docProps/core.xml regardless.

    docx has no unambiguous discrete unit (paragraphs aren't stable locators), so
    units=[]/unit_key=None — but `tier` is recorded so the provider can mark
    structured output (R4)."""
    text = _docx_via_python_docx(filepath)
    tier = 'python-docx' if text else ''
    if not text:
        text = _pandoc_to_markdown(filepath, 'docx')
        tier = 'pandoc' if text else ''
    if not text:
        text = _zip_text_from(filepath, ['word/document.xml'], 't')
        tier = 'zipfile' if text else ''
    title = _office_core_title(filepath)
    return {'title': title, 'text': text[:50000],
            'tier': tier, 'units': [], 'unit_key': None}


def _extract_xlsx(filepath):
    """Excel: pandoc's xlsx reader emits each sheet as a markdown
    section with tables — much more useful than a flat string of cell
    values. Stdlib fallback preserves the old behavior so ingest
    works without pandoc.

    Per-sheet units (R5) come from the structured (openpyxl) tier when it ran,
    else the stdlib helper — so locators survive every tier, keyed by sheet name."""
    text, units = _xlsx_via_openpyxl(filepath)
    tier = 'openpyxl' if text else ''
    if not text:
        text = _pandoc_to_markdown(filepath, 'xlsx')
        tier = 'pandoc' if text else ''
    if not text:
        try:
            with zipfile.ZipFile(filepath) as zf:
                sheet_entries = sorted(
                    (n for n in zf.namelist()
                     if n.startswith('xl/worksheets/sheet') and n.endswith('.xml')),
                    key=_entry_index,
                )
        except (zipfile.BadZipFile, OSError):
            sheet_entries = []
        shared_text = _zip_text_from(filepath, ['xl/sharedStrings.xml'], 't')
        inline_text = _zip_text_from(filepath, sheet_entries, 't')
        text = (shared_text + '\n' + inline_text).strip()
        tier = 'zipfile' if text else ''
    if not units:
        units = _xlsx_units(filepath)
    title = _office_core_title(filepath)
    return {'title': title, 'text': text[:50000],
            'tier': tier, 'units': units, 'unit_key': 'sheet'}


def _extract_pptx(filepath):
    """PowerPoint: pandoc emits each slide as a level-1 heading with
    its text under it (bullet points preserved). The stdlib fallback
    produces flat concatenated text runs.

    Per-slide units (R5) come from the structured (python-pptx) tier when it ran,
    else the stdlib helper — keyed by 1-indexed slide number."""
    text, units = _pptx_via_python_pptx(filepath)
    tier = 'python-pptx' if text else ''
    if not text:
        text = _pandoc_to_markdown(filepath, 'pptx')
        tier = 'pandoc' if text else ''
    if not text:
        try:
            with zipfile.ZipFile(filepath) as zf:
                slide_entries = sorted(
                    (n for n in zf.namelist()
                     if n.startswith('ppt/slides/slide') and n.endswith('.xml')),
                    key=_entry_index,
                )
        except (zipfile.BadZipFile, OSError):
            slide_entries = []
        text = _zip_text_from(filepath, slide_entries, 't')
        tier = 'zipfile' if text else ''
    if not units:
        units = _pptx_units(filepath)
    title = _office_core_title(filepath)
    return {'title': title, 'text': text[:50000],
            'tier': tier, 'units': units, 'unit_key': 'slide'}


def _extract_epub(filepath):
    """EPUB: pandoc handles spine order, chapter headings, and inline
    formatting correctly — much better than our tag-stripping fallback.
    The fallback walks .xhtml/.html files in name order (close-enough
    reading order) and strips HTML tags via regex.

    EPUB chapter files aren't a stable, user-meaningful locator unit, so
    units=[]/unit_key=None — but `tier` is recorded for the structured
    marker (R4)."""
    text = _pandoc_to_markdown(filepath, 'epub')
    if text:
        return {'title': '', 'text': text[:50000],
                'tier': 'pandoc', 'units': [], 'unit_key': None}

    try:
        with zipfile.ZipFile(filepath) as zf:
            doc_entries = sorted(
                n for n in zf.namelist()
                if (n.endswith('.xhtml') or n.endswith('.html')) and not n.startswith('META-INF/')
            )
            chunks = []
            for entry in doc_entries:
                try:
                    raw = zf.read(entry).decode('utf-8', errors='replace')
                except (OSError, KeyError):
                    continue
                stripped = re.sub(r'<[^>]+>', ' ', raw)
                stripped = re.sub(r'\s+', ' ', stripped).strip()
                chunks.append(stripped)
    except (zipfile.BadZipFile, OSError):
        chunks = []
    text = '\n\n'.join(chunks)[:50000]
    return {'title': '', 'text': text,
            'tier': 'zipfile' if text else '', 'units': [], 'unit_key': None}


def _office_core_title(filepath):
    """Pull <dc:title> from docProps/core.xml — the common metadata
    file across all Office 2007+ formats. Returns '' if absent."""
    try:
        with zipfile.ZipFile(filepath) as zf:
            if 'docProps/core.xml' not in zf.namelist():
                return ''
            raw = zf.read('docProps/core.xml')
            root = ET.fromstring(raw)
            for el in root.iter():
                tag = el.tag.rsplit('}', 1)[-1] if '}' in el.tag else el.tag
                if tag == 'title' and el.text:
                    return el.text.strip()
    except (zipfile.BadZipFile, ET.ParseError, OSError):
        pass
    return ''

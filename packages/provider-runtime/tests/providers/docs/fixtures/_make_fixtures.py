"""Generate tiny test fixtures for DocsProvider. Run from repo root:

    uv run --extra office python packages/provider-runtime/tests/providers/docs/fixtures/_make_fixtures.py

Produces small.docx, small.xlsx, small.pptx, small.epub in this directory.
Each fixture has a known title + body string the tests assert against.
"""

from __future__ import annotations

import sys
import zipfile
from pathlib import Path

OUT = Path(__file__).parent
KNOWN_TITLE = "Test Document"
KNOWN_BODY = "Body text for DocsProvider testing."


def make_docx() -> None:
    from docx import Document

    doc = Document()
    doc.core_properties.title = KNOWN_TITLE
    doc.add_heading(KNOWN_TITLE, level=1)
    doc.add_paragraph(KNOWN_BODY)
    doc.save(OUT / "small.docx")
    print(f"wrote {OUT / 'small.docx'}")


def make_xlsx() -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Test Sheet"
    ws["A1"] = "Header"
    ws["A2"] = KNOWN_BODY
    wb.save(OUT / "small.xlsx")
    print(f"wrote {OUT / 'small.xlsx'}")


def make_pptx() -> None:
    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = KNOWN_TITLE
    slide.placeholders[1].text = KNOWN_BODY
    prs.save(OUT / "small.pptx")
    print(f"wrote {OUT / 'small.pptx'}")


def make_epub() -> None:
    """Minimal valid EPUB 3 zip: mimetype + META-INF/container.xml +
    OEBPS/content.opf + OEBPS/chapter1.xhtml.

    Hand-rolled because ebooklib isn't a dep. pandoc reads this fine."""
    out_path = OUT / "small.epub"

    container_xml = """<?xml version="1.0" encoding="UTF-8"?>
<container version="1.0" xmlns="urn:oasis:names:tc:opendocument:xmlns:container">
  <rootfiles>
    <rootfile full-path="OEBPS/content.opf" media-type="application/oebps-package+xml"/>
  </rootfiles>
</container>
"""

    content_opf = f"""<?xml version="1.0" encoding="UTF-8"?>
<package xmlns="http://www.idpf.org/2007/opf" version="3.0" unique-identifier="bookid">
  <metadata xmlns:dc="http://purl.org/dc/elements/1.1/">
    <dc:identifier id="bookid">urn:uuid:arcus-test-12345</dc:identifier>
    <dc:title>{KNOWN_TITLE}</dc:title>
    <dc:language>en</dc:language>
  </metadata>
  <manifest>
    <item id="ch1" href="chapter1.xhtml" media-type="application/xhtml+xml"/>
  </manifest>
  <spine>
    <itemref idref="ch1"/>
  </spine>
</package>
"""

    chapter1 = f"""<?xml version="1.0" encoding="UTF-8"?>
<!DOCTYPE html>
<html xmlns="http://www.w3.org/1999/xhtml">
<head><title>{KNOWN_TITLE}</title></head>
<body>
  <h1>{KNOWN_TITLE}</h1>
  <p>{KNOWN_BODY}</p>
</body>
</html>
"""

    # mimetype must be the first entry, uncompressed, exactly "application/epub+zip"
    with zipfile.ZipFile(out_path, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr(zipfile.ZipInfo("mimetype"), "application/epub+zip", zipfile.ZIP_STORED)
        zf.writestr("META-INF/container.xml", container_xml)
        zf.writestr("OEBPS/content.opf", content_opf)
        zf.writestr("OEBPS/chapter1.xhtml", chapter1)

    print(f"wrote {out_path}")


def main() -> int:
    make_docx()
    make_xlsx()
    make_pptx()
    make_epub()
    return 0


if __name__ == "__main__":
    sys.exit(main())
